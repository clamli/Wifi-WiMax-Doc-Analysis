# -*- coding: utf-8 -*-
"""
Created on Mon Sep 16 23:55:37 2018

@author: zhenkai wang(Kay)
parallel version of word-freq-matrix building.
"""
import fitz, re, os, pickle, multiprocessing
#import nltk
from nltk import word_tokenize, pos_tag, ne_chunk, chunk
from nltk.corpus import treebank
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import string
from sklearn.feature_extraction import DictVectorizer
from sklearn.feature_extraction.text import CountVectorizer
from functools import partial
import warnings

pdfDir11_1 = "F:\\WireLessNLPGRA\\gitFiles\\Wifi-WiMax-Doc-Analysis\\Spider_Kay\\80211grouper"
pdfDir11_2 = "F:\\WireLessNLPGRA\\gitFiles\\Wifi-WiMax-Doc-Analysis\\Spider_Kay\\80211mentor"
pdfDir16_1 = "F:\\WireLessNLPGRA\\gitFiles\\Wifi-WiMax-Doc-Analysis\\Spider_Kay\\80216mentor"


#path1 = "F:\\WireLessNLPGRA\\code\\NamedEntityRecognition\\some_proposal\\C802162a-01_01.pdf"
#path2 = "F:\\WireLessNLPGRA\\code\\NamedEntityRecognition\\some_proposal\\C802162a-01_06.pdf"
techTermPath = "words"
abbListPath = "Abbreviation_from_files_1000_by.dat"
dictPath = "dict"
fileNamesPath = ""
abbWLDictPath = "80211grouper_AB_1393file_wl_dict.dat"
#pdfDir =  "F:\\WireLessNLPGRA\\80216CompProj"#change this to the pdf folders
numCount = 0
setting = {"80211mentor":{
        "fileName":"80211mentorIncludedFileNames.dat",
        "abbWLDict": "80211mentorIncludedAbbWLDict.dat",
        "fileDir": "F:\\WireLessNLPGRA\\gitFiles\\Wifi-WiMax-Doc-Analysis\\Spider_Kay\\80211mentor"
        }
            
    }

pdfDir =  pdfDir11_1
from fileProcessing import *
from techTermProcessing import *

def getPage(path):
# =============================================================================
#     get all the pages of file, return a list of string, each string is content of one page
# =============================================================================
    doc = fitz.open(path)
    content = []
    for i in range(doc.pageCount):
        content.append(doc.getPageText(i))
    return content
def abstractExtracter(firstPage):
# =============================================================================
#     extract abstract of the first page
# =============================================================================
    pAbstract = r'^Abstract\s(.+)'
    candidate = [re.findall(pAbstract,line) for line in firstPage.split('\n')]
    abstract = []
    for sen in candidate:
        if (len(sen) > 0):
            abstract.append(sen[0])
    return abstract[0]
def getStopWords():
# =============================================================================
#     #we return the stopWords in NLTK database
# =============================================================================
    return set(stopwords.words('english'))
def transform(word):
# =============================================================================
#     if a word is all upper case then return itself otherwise should return word lower case 
# =============================================================================
    isAllUpper = (word == word.upper())
    if isAllUpper:
        pass
    else:
        word = word.lower()
    return word

from nltk.corpus import words
english = words.words()
english = [word.lower() for word in english]
def getTechTerms(words, oneGram, twoGram, threeGram, extractedAbb, abbSet, pattern = "AB"):
    global english
# =============================================================================
#     #given a list of words from a doc, and a list of tecnical terms, find all occurance of technical word
# =============================================================================
#    p1Match = r'^([A-Z]{2,})$'
#    p2Contain = r'802'
    print("pattern\t", pattern)
    result = []
    wordNum = len(words)
    if ("tech" in pattern):        
        for i in range(wordNum):
            word = words[i]

            #oneGram
            if (word in oneGram):
                result.append(word)                
            #twoGram
            if (word in twoGram and i < wordNum - 1 and (word + " " + words[i + 1]) == twoGram[word]):
                result.append(twoGram[word])               
                    
            #threeGram
            if (word in threeGram and i < wordNum - 2 and (word + " " + words[i + 1] + " " + words[i + 2]) == threeGram[word]):
                result.append(threeGram[word])

    #delete stopword
    stopWords = getStopWords()
    wordsFiltered = []
    for w in result:
        if w not in stopWords:
            wordsFiltered.append(w)
    appendList = []
    if ("AB" in pattern):
        if (len(extractedAbb) > 0):
            for word in extractedAbb:
                #prevent double adding abb
                if not (word in wordsFiltered):
                    appendList.append(word)#should not use wordsFilteredList.append, or it will only record one
#        else:
#            setFilterd = set(wordsFiltered)
#            for word in words:
#                if (word not in setFilterd and word.isalpha() and len(word)> 1):
#                    if not isAbbrev(word) and word.lower() not in english:
#                        appendList.append(word)
    finalResult = wordsFiltered + appendList
    print("num of technical terms got\t", len(finalResult))
    return finalResult

def moreThanOneUpper(word):
    accum = 0
    for ch in word:
        if ch.isupper():
            accum += 1
            if accum > 1:
                return True
    return False

def isAbbrev(word):
    if word.isupper() or moreThanOneUpper(word):
        return True
    return False

def notInWords(word):
    for pott in words_tec:
        if word in pott:
            return False
    return True
#def isTechTerms(word, techTermSet):
#    return word
##    for term in listOfTechTerm:
##                #correct match
##                if word == term:
##                    isTerm = True
##                splited_terms = term.split()
##                
##                #see if word is the phrase:
##                if len(splited_terms) > 1:
##                    isEqual = True
##                    for idx in range(len(splited_terms)):
##                        if (i + idx > wordNum - 1):
##                            break
##                        if (words[i + idx] != splited_terms[idx]):
##                            isEqual = False
##                    if (isEqual):
##                        result.append(term)
#def isAbbreviation(word, abbSet):
##    return (word.upper() == word)#then all words that are numbers, will return true
#    #1.len > 1;
#    #2.num of upper case characters >= 1/2 len
#    numChar = len(word)
#    if (numChar <= 1):
#        return False
#    acc = 0
#    for c in word:
#        if (c.isupper()):
#            acc+=1
#    if (acc >= numChar / 2):
#        return True
#    return False
    
def hasCharDig(word):
# =============================================================================
#     #check whether a word has both digit and char
# =============================================================================

    hasChar = False
    hasWord = False
    for i in range(len(word)):
        c = word[i]
        if(c.isdigit()):
            hasChar = True
        if(c.isalpha()):
            hasWord = True
    return hasChar & hasWord
def getListOfTechTerms(dictAbb, *paths, isReset = False):
# =============================================================================
#     #functions: from both <abb, full title> dictionary and tech terms list, extract
#     #add both key and value in dictionary, add all the tech terms in tech terms list
#     #isReset: if true, then read all the files, if false, just try to read the dumped record of tech terms
# 
# =============================================================================
    try:
        if (isReset):
            print("reseting all technical terms")
            raise Exception("reseting")
        with open("techTerms.dat", "rb") as f:
            technicalWords = pickle.load(f)
    except:
        technicalWords = []
        for path in paths:
            with open (path, 'rb') as fp:
                technicalWords += pickle.load(fp)
        #Strip
        for i in range(len(technicalWords)):
            technicalWords[i] = technicalWords[i].strip()
        
        for k, v in dictAbb.items():
            technicalWords.append(k)
            technicalWords.append(v)
        temp = list(set(technicalWords))
        technicalWords = []
        for word in temp:
            if (len(word) > 1 and word.lower() != "ieee"):
                
                technicalWords.append(word)
#        print("finish newly loading technical terms list, with num of words:" + str(len(technicalWords)))
        with open("techTerms.dat", "wb") as f:
            pickle.dump(technicalWords, f)
            
    
#    print("finish loading technical terms list, with num of words:" + str(len(technicalWords)))
        
    return technicalWords

from pptx import Presentation
import docx2txt
def getContent(textPath):
# =============================================================================
# #    input is the path of pdf file, return all of the words(segemented by \s \n)
# #    remove all punctuation
# =============================================================================
    words = []
    if (textPath.endswith(".pdf")):
        
        content = getPage(textPath)
        text1 = '.'.join(content)
    
        #words = nltk.word_tokenize(text1)
        sents = sent_tokenize(text1)
        #delete punctuation
        for sent in sents:
            ws = [word.strip(string.punctuation) for word in re.split(r'\s|\n', sent)]
            ws = [w for w in ws if len(w) > 1]
            words += ws
    if (textPath.endswith(".doc") or textPath.endswith(".docx")):
        text = docx2txt.process(textPath)
        sents = re.split("\\n+|\\t+|\.+", text)
        for sent in sents:
            if (sent is None ): continue
            ws =  [word.strip(string.punctuation) for word in re.split(r'\s|\n', sent)]
            ws = [w for w in ws if len(w) > 1]
            words+=ws
    if (textPath.endswith(".ppt") or textPath.endswith(".pptx")):
        #python pptx  is good at extracting content in pptx, but ppt is not ok
        sents = []
        prs = Presentation(textPath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        sents.append(run.text)
        for sent in sents:
            ws =  [word.strip(string.punctuation) for word in re.split(r'\s|\n', sent)]
            ws = [w for w in ws if len(w) > 1]
            words += ws
    #delete stopword
#     stopWords = set(stopwords.words('english'))
#     wordsFiltered = []
#     for w in words:
#        if w not in stopWords:
#            wordsFiltered.append(w)

    wordsFiltered = [lowerlize(word) for word in words]
    wordsFiltered = [word for word in wordsFiltered if word is not None]
    return wordsFiltered
def lowerlize(word):
    try:
            
        if (word[0].isupper() and word != word.upper()):
            return word.lower()
        else:
            return word
    except:
        return None



def getList(file, fileDir, oneGram, twoGram, threeGram, abb_all, abbSet, dictAbb, pattern = "AB"):
# =============================================================================
#     #get a list of word lists, each element is the word extracted from file from files
#     #replace all the words  that occurring in key set of dictAbb by corresponding full titles
# =============================================================================
    global numCount

    
    wordsFiltered = getContent(os.path.join(fileDir, file))
    techTerms = getTechTerms(wordsFiltered, oneGram, twoGram, threeGram, abb_all.get(file.split("\\")[-1], []), abbSet, pattern)
#    for i in range(len(techTerms)):
#        abb = techTerms[i]
#        fullTitle = dictAbb.get(abb, None)
#        if fullTitle is not None:
#            techTerms[i] = fullTitle
##                    print("changeing ", abb, " \t to ", fullTitle)
#        listOfWords.append(techTerms)
#    print(numCount, "\t", file)
#    numCount += 1
#    try:
#        wordsFiltered = getContent(file)
#        techTerms = getTechTerms(wordsFiltered, listOfTechTerms, pattern)
#        for i in range(len(techTerms)):
#            abb = techTerms[i]
#            fullTitle = dictAbb.get(abb, None)
#            if fullTitle is not None:
#                techTerms[i] = fullTitle
##                    print("changeing ", abb, " \t to ", fullTitle)
#            listOfWords.append(techTerms)
#        print(numCount, "\t", file)
#        numCount += 1
#            #here delete the duplicate one
#    except:
#        pass#if can't open the file then do nothing
    return (techTerms, file.split("\\")[-1])


def getDictOfAbb(*paths):
# =============================================================================
#     #this file is a list of list, each of the sublist contains two string, first one is abbreviation, second one is full title
#     #make it into a dictionary
# =============================================================================
    dictAbb = {}
    for path in paths:
        with open(path, 'rb') as fp:
            l = pickle.load(fp)
#     for key, value in dictAbb.iteritems():
    for i in range(len(l)):
        dictAbb[l[i][0]] = l[i][1].lower()

    return dictAbb
from itertools import product
if __name__ == "__main__":
    warnings.filterwarnings("ignore")
    #key paras
#    minCount = 1
#    numStart = 0
#    numEnd = 2000
#    prefix = "AB"
    prefix = "tech_AB"
#    prefix = "tech"    
    if (prefix == "AB_tech"):
        print("extracting both technical terms and abbreviation")
    elif(prefix == "tech"):
        print("extracting techincal terms")
    else:
        print("extractng abbreviation")
    #prepocessing
    dictAbb1 = getDictOfAbb(dictPath)
    #if put trus as third para, then it will reextract the tech terms from files
    listOfTechTerms1 = getListOfTechTerms(dictAbb1, techTermPath, True)#true means resetting
    techTermSet1 = set(listOfTechTerms1)
    
    dictOfAbb = pickle.load(open(abbListPath, "rb"))
    listOfAbb = []
    for k, v in dictOfAbb.items():
        listOfAbb = listOfAbb + v
    
    abbSet1 = set(listOfAbb)
    print("finish loading")
    #get all the files
    
    #loading all the paras for getListfunction
#    fileNames = getSelectedFileNames(pdfDir, numStart, numEnd, pdfDir.split("\\")[-1] + "randomIdx.dat");
    fileNames = pickle.load(open(fileNamesPath, "rb"))
    oneGram = pickle.load(open("techTrem_oneGram_set.dat", "rb"))
    twoGram = pickle.load(open("techTerm_twoGram_dict.dat", "rb"))
    threeGram = pickle.load(open("techTerm_threeGram_dict.dat", "rb"))
    abb_all = pickle.load(open(abbWLDictPath, "rb"))
    
    
    print("start parallal processing")
    np = 14
    useAbbreviation = True#true means use both abbreviation and tech terms
    p = multiprocessing.Pool(np)
    output = p.map(partial(getList, pdfDir = pdfDir, oneGram = oneGram, twoGram = twoGram, threeGram = threeGram, abbSet = abbSet1, abb_all = abb_all, dictAbb = dictAbb1, pattern = prefix),  [file for file in fileNames])
#    print(len(output[0]))
    listOfWordDict = {}
    computedFiles = []
    for out in output:
        if (len(out[0]) > 0):
            listOfWordDict[out[1]] = out[0]#each element in listOfWordLists is a tuple, whose first elem is fileName, second is the list of words
#            computedFiles.append(out[1])
#    print(listOfWordLists[0])
    print("number of files:", len(computedFiles))
    
    #calculate the word-freq-matrix
#    vec = CountVectorizer(min_df = minCount,tokenizer=lambda doc: doc, lowercase=False)
#    wordFreqArray = vec.fit_transform(listOfWordLists)
#    featureName = vec.get_feature_names()
    
    #save key variables

    keyWord = pdfDir.split("\\")[-1] +  "_" + prefix +"_" + str(len(listOfWordDict))+ "file_"
    pickle.dump(listOfWordDict, open(keyWord + "wl_dict.dat", "wb"))
    pickle.dump(listOfWordDict, open("wl_dict.dat", "wb"))
